from docx2pdf import convert
from fpdf import FPDF
from pdf2docx import Converter
from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image
from PyPDF2 import PdfMerger
import pandas as pd
import img2pdf
import pdfkit
import tabula
import os


def add_ext(filename, ext):
    return filename if filename.endswith(ext) else filename + ext


def get_valid_input(prompt, valid_options):
    while True:
        choice = input(prompt)
        if choice in valid_options:
            return choice
        print(f'\nInvalid option. Please select from: {", ".join(valid_options)}')


def get_valid_file_path(prompt):
    while True:
        path = input(prompt)
        if os.path.exists(path):
            return path
        print('\nFile not found. Please enter a valid file path.')


# main menu
print('[1] - Turn a file into a PDF\n'
      '[2] - Turn PDF to another file\n'
      '[3] - Merge two PDFs\n'
      '[4] - Turn an image file into another image file type')

main_choice = get_valid_input('What would you like to do: ', ['1', '2', '3', '4'])

# convert to pdf
if main_choice == '1':
    print('[1] - .docx\n'
          '[2] - .txt\n'
          '[3] - .jpeg / .jpg\n'
          '[4] - .png\n'
          '[5] - .webp\n'
          '[6] - .csv\n'
          '[7] - .xlsx\n'
          '[8] - .pptx')
    
    format_choice = get_valid_input('What would you like to do: ', 
                                  ['1', '2', '3', '4', '5', '6', '7', '8'])
    
    inp_file_path = get_valid_file_path('File path for the file you want changed: ')
    opt_file_name = add_ext(input('What do you want to name the file: '), '.pdf')

    # docx to pdf
    if format_choice == '1':
        convert(inp_file_path, opt_file_name)
    
    # txt to pdf
    elif format_choice == '2':
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        with open(inp_file_path, "r") as f:
            for line in f:
                pdf.cell(200, 10, txt=line, ln=1, align="L")
        pdf.output(opt_file_name)
    
    # image to pdf
    elif format_choice in ['3', '4', '5']:
        with open(opt_file_name, "wb") as f:
            f.write(img2pdf.convert(inp_file_path))
    
    # csv/xlsx to pdf
    elif format_choice in ['6', '7']:
        df = pd.read_csv(inp_file_path) if format_choice == '6' else pd.read_excel(inp_file_path)
        pdfkit.from_string(df.to_html(), opt_file_name)
    
    # pptx to pdf
    elif format_choice == '8':
        prs = Presentation(inp_file_path)
        temp_dir = "temp_slides"
        os.makedirs(temp_dir, exist_ok=True)
        
        for i, slide in enumerate(prs.slides):
            slide.export(f"{temp_dir}/slide_{i}.png", "PNG")
        
        images = convert_from_path(f"{temp_dir}/*.png")
        images[0].save(opt_file_name, "PDF", save_all=True, append_images=images[1:])
        
        # cleanup
        for file in os.listdir(temp_dir):
            os.remove(os.path.join(temp_dir, file))
        os.rmdir(temp_dir)

# convert from pdf
elif main_choice == '2':
    print('[1] - .docx\n'
          '[2] - .jpg\n'
          '[3] - .jpeg\n'
          '[4] - .png\n'
          '[5] - .webp\n'
          '[6] - .xlsx\n'
          '[7] - .pptx')
    
    format_choice = get_valid_input('What would you like to do: ',
                                  ['1', '2', '3', '4', '5', '6', '7'])
    
    inp_file_path = get_valid_file_path('File path for the PDF: ')
    
    # pdf to docx
    if format_choice == '1':
        opt_file_name = add_ext(input('Output filename: '), '.docx')
        cv = Converter(inp_file_path)
        cv.convert(opt_file_name)
        cv.close()
    
    # pdf to image
    elif format_choice in ['2', '3', '4', '5']:
        ext_map = {'2': '.jpg', '3': '.jpeg', '4': '.png', '5': '.webp'}
        opt_file_name = add_ext(input('Output filename: '), ext_map[format_choice])
        images = convert_from_path(inp_file_path)
        images[0].save(opt_file_name)
    
    # pdf to xlsx
    elif format_choice == '6':
        opt_file_name = add_ext(input('Output filename: '), '.xlsx')
        pd.concat(tabula.read_pdf(inp_file_path, pages='all')).to_excel(opt_file_name, index=False)
    
    # pdf to pptx
    elif format_choice == '7':
        opt_file_name = add_ext(input('Output filename: '), '.pptx')
        temp_dir = "temp_slides"
        os.makedirs(temp_dir, exist_ok=True)
        
        images = convert_from_path(inp_file_path)
        image_paths = []
        for i, image in enumerate(images):
            img_path = f"{temp_dir}/slide_{i}.png"
            image.save(img_path, "PNG")
            image_paths.append(img_path)
        
        prs = Presentation()
        for img_path in image_paths:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
        
        prs.save(opt_file_name)
        
        # cleanup
        for file in image_paths:
            os.remove(file)
        os.rmdir(temp_dir)

# merge pdfs
elif main_choice == '3':
    inp_file_path1 = get_valid_file_path('Enter the path of first PDF file: ')
    inp_file_path2 = get_valid_file_path('Enter the path of second PDF file: ')
    opt_file_name = add_ext(input('Enter the output file name: '), '.pdf')
    
    merger = PdfMerger()
    merger.append(inp_file_path1)
    merger.append(inp_file_path2)
    merger.write(opt_file_name)
    merger.close()
    
    print(f'\nPDFs merged successfully! Output saved as: {opt_file_name}')

# convert image formats
else:
    inp_file_path = get_valid_file_path('Enter the path of the image file: ')
    
    valid_formats = ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff']
    print('\nSupported image formats:', ', '.join(valid_formats))
    
    while True:
        out_ext = input('Enter the desired output image format (e.g. .png): ').lower()
        if out_ext in valid_formats:
            break
        print('\nUnsupported format. Please enter a valid image extension.')
    
    opt_file_name = add_ext(input('Enter the output file name: '), out_ext)
    
    img = Image.open(inp_file_path)
    if out_ext in ['.jpg', '.jpeg'] and img.mode == 'RGBA':
        img = img.convert('RGB')
    
    img.save(opt_file_name)
    print(f'\nImage converted successfully! Output saved as: {opt_file_name}')